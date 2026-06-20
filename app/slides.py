"""
スライド生成サービス層

方針:
1. Markdown の見出し1/2/3ごとに1スライドへ分解（決定論パーサ・LLM不要）
   - h4以降・箇条書き・段落は直近スライドの本文に束ねる
2. 各スライドを GPT Image 2.0 で1枚絵として生成
   - プロジェクトの見本画像があれば images.edit の参照に使う（無ければ images.generate）
   - プロジェクトごとの生成プロンプトをスタイル指示に流用
   - 任意で「修正指示(fix)」をプロンプト末尾に追加して再生成
3. PNG 群を python-pptx で「1枚画像=1スライド（16:9 全面）」の PPTX に組立

出力先は OUTPUT_ROOT/<プロジェクト名>/ （既定は iCloud Drive の Downloads）。
"""
from __future__ import annotations

import os
import io
import re
import json
import base64
from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches

# ---------------------------------------------------------------------------
# パス設定
# ---------------------------------------------------------------------------

OUTPUT_ROOT = Path(
    os.environ.get(
        "OUTPUT_ROOT",
        Path.home() / "Library/Mobile Documents/com~apple~CloudDocs/Downloads",
    )
)

SLIDE_IMAGE_MODEL = os.environ.get("SLIDE_IMAGE_MODEL", "gpt-image-2")

# gpt-image の landscape サイズ（3:2）。生成後に 16:9 へ縦圧縮フィットする。
GEN_SIZE = "1536x1024"
TARGET_W, TARGET_H = 1536, 864  # 16:9

# プロジェクト内の固定ファイル名
MARKDOWN_FILE = "markdown.md"
PROMPT_FILE = "_prompt.txt"
REFERENCE_FILE = "_reference.png"
PLAN_FILE = "plan.json"
CONFIG_FILE = "_config.json"


# ---------------------------------------------------------------------------
# プロンプトテンプレート（新規作成時に選択・編集できる雛形）
# ---------------------------------------------------------------------------

PROMPT_TEMPLATES: dict[str, str] = {
    "SHI-GYO": (
        "あなたは最高のプレゼンテーションスライドメーカーです。\n"
        "参照画像のスライドのテイスト（デザイン・色味）を踏襲して、"
        "情報発信プロジェクト「Dr.SHI-GYO」に関連するスライドを生成してください。\n"
        "・文字数はできるだけ減らす（キーワード中心）\n"
        "・高齢者も見るセミナーなので小さい文字は禁止（引用文献情報のみ小さく）\n"
        "・少ない文字数・大きな文字で、イラスト・図表・写真も使って直感的に理解させる\n"
        "【右下にロゴ用スペースを空けるルール・厳守】\n"
        "・右下の角に、ロゴやワイプを置くための小さな空きスペース"
        "（背景のまま・文字も図版も置かない／目安は横幅の約15%・高さの約20%）を必ず確保する。\n"
        "・カードや図版・本文は、その右下の角を避けて配置する。"
        "右下にカードが重なる場合は、最後の項目を左に寄せる・カードを一段上げる・"
        "下段を3列ではなく左寄せにするなどして、右下の角だけは背景を残すこと。\n"
        "・ただし空けるのは右下の角だけ。右端を縦一列にまるごと空けたり、"
        "画面の右3〜4割を縦長の空白にするのは絶対に禁止。\n"
        "・左・中央・上部、そして右上〜右中央は画面いっぱいに使い、横幅を端近くまで使い切ること。"
    ),
    "シンプル": (
        "クリーンでミニマルなビジネスプレゼンのスライド。"
        "白基調の背景、余白を広めに取り、見出しを大きく、本文は読みやすい階層で配置。"
        "アクセントカラーを1色だけ使う。装飾は控えめに。"
    ),
    "ポップ": (
        "明るくカラフルでポップなプレゼンスライド。"
        "親しみやすい配色、丸みのある図形、わかりやすいアイコン的イラストを添える。"
        "見出しは目を引く大きさで。"
    ),
    "ダーク": (
        "ダークテーマのモダンなプレゼンスライド。"
        "濃紺〜黒の背景に明るい文字、ネオン的なアクセント1色。"
        "テック系の洗練された印象。"
    ),
    "和風・落ち着き": (
        "落ち着いた和モダンのプレゼンスライド。"
        "生成り・墨・若干の差し色。余白を活かし、上品で読みやすいレイアウト。"
    ),
}

DEFAULT_PROMPT = PROMPT_TEMPLATES["SHI-GYO"]

# テンプレートに同梱する「デフォルト見本デザイン画像」。
# キーは PROMPT_TEMPLATES のテンプレート名、値は app/template_assets/ 内のファイル名。
# 新規作成時に見本画像が未アップロードなら、ここで指定した画像が初期見本になる。
# （プロジェクトで見本画像をアップロードすると、そのプロジェクトの見本に差し替わる）
TEMPLATE_ASSETS_DIR = Path(__file__).resolve().parent / "template_assets"

TEMPLATE_REFERENCES: dict[str, str] = {
    "SHI-GYO": "shigyo_reference.png",
}

# テンプレートごとの「右下ネガティブスペース（ロゴ/ワイプ用の余白）」設定。
# プロンプトで余白を空けさせたうえで、生成後に下記の矩形を背景色で軽く均して
# きれいな余白に仕上げる（端のフレームは inset 分だけ残す）。
#   right / bottom : 画像の右端・下端からの割合（空ける領域の幅・高さ）
#   inset          : 外周フレームを残すための内側マージン（幅に対する割合）
#   feather        : 合成マスクのぼかし量（幅に対する割合）
TEMPLATE_NEGATIVE_SPACE: dict[str, dict] = {
    "SHI-GYO": {"right": 0.15, "bottom": 0.20, "inset": 0.014, "feather": 0.018},
}


def template_reference_path(template_name: str) -> Path | None:
    """テンプレートに同梱されたデフォルト見本画像のパスを返す（無ければ None）。"""
    fname = TEMPLATE_REFERENCES.get(template_name or "")
    if not fname:
        return None
    p = TEMPLATE_ASSETS_DIR / fname
    return p if p.exists() else None


def templates_with_reference() -> list[str]:
    """デフォルト見本画像を持つテンプレート名の一覧。"""
    return [name for name in TEMPLATE_REFERENCES if template_reference_path(name)]


def template_negative_space(template_name: str) -> dict | None:
    """テンプレートに定義された右下ネガティブスペース設定を返す（無ければ None）。"""
    ns = TEMPLATE_NEGATIVE_SPACE.get(template_name or "")
    return dict(ns) if ns else None


# ---------------------------------------------------------------------------
# プロジェクトのパス
# ---------------------------------------------------------------------------

def project_dir(project_name: str) -> Path:
    d = OUTPUT_ROOT / project_name
    d.mkdir(parents=True, exist_ok=True)
    return d


def markdown_path(project_name: str) -> Path:
    return project_dir(project_name) / MARKDOWN_FILE


def prompt_path(project_name: str) -> Path:
    return project_dir(project_name) / PROMPT_FILE


def reference_path(project_name: str) -> Path:
    return project_dir(project_name) / REFERENCE_FILE


def apply_template_reference(project_name: str, template_name: str) -> bool:
    """テンプレートのデフォルト見本画像をプロジェクトの _reference.png に設定する。

    既にプロジェクト独自の見本画像がある場合は上書きしない。
    コピーした場合は True、対象が無い／既存ありの場合は False を返す。
    """
    src = template_reference_path(template_name)
    if src is None:
        return False
    dst = reference_path(project_name)
    if dst.exists():
        return False
    import shutil

    project_dir(project_name)
    shutil.copyfile(src, dst)
    return True


def config_path(project_name: str) -> Path:
    return project_dir(project_name) / CONFIG_FILE


def read_project_config(project_name: str) -> dict:
    p = config_path(project_name)
    if not p.exists():
        return {}
    try:
        return json.loads(p.read_text(encoding="utf-8"))
    except Exception:
        return {}


def write_project_config(project_name: str, config: dict) -> None:
    config_path(project_name).write_text(
        json.dumps(config, ensure_ascii=False, indent=2), encoding="utf-8"
    )


def init_project_from_template(project_name: str, template_name: str) -> dict:
    """テンプレート選択時の初期化（見本画像コピー＋設定保存）をまとめて行う。

    返り値: {"applied_reference": bool, "negative_space": dict | None}
    """
    applied_reference = apply_template_reference(project_name, template_name)
    ns = template_negative_space(template_name)
    if ns:
        cfg = read_project_config(project_name)
        cfg["negative_space"] = ns
        write_project_config(project_name, cfg)
    return {"applied_reference": applied_reference, "negative_space": ns}


def plan_path(project_name: str) -> Path:
    return project_dir(project_name) / PLAN_FILE


def slide_png_path(project_name: str, index: int) -> Path:
    return project_dir(project_name) / f"slide_{index:02d}.png"


def pptx_path(project_name: str) -> Path:
    return project_dir(project_name) / f"{project_name}.pptx"


def read_markdown(project_name: str) -> str:
    p = markdown_path(project_name)
    return p.read_text(encoding="utf-8") if p.exists() else ""


def read_prompt(project_name: str) -> str:
    p = prompt_path(project_name)
    return p.read_text(encoding="utf-8") if p.exists() else DEFAULT_PROMPT


def _client() -> OpenAI:
    return OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))


# ---------------------------------------------------------------------------
# ① Markdown → スライド構成（見出し1/2/3ごとに1枚）
# ---------------------------------------------------------------------------

_HEADING_RE = re.compile(r"^(#{1,3})\s+(.*)$")


def parse_markdown_to_slides(md: str) -> list[dict]:
    """Markdown を見出し1/2/3ごとに1スライドへ分解する。

    - `#` `##` `###` が現れるたびに新スライド境界
    - h4以降・箇条書き・番号付き・段落は直近スライドの本文(lines)に束ねる
    - 最初の見出しより前に本文がある場合は、無題スライドとして拾う
    """
    slides: list[dict] = []
    current: dict | None = None

    for raw in (md or "").splitlines():
        line = raw.rstrip()
        m = _HEADING_RE.match(line)
        if m:
            if current is not None:
                slides.append(current)
            current = {"title": m.group(2).strip(), "lines": [], "illustration": ""}
            continue

        stripped = line.strip()
        if not stripped:
            continue

        if current is None:
            # 最初の見出し前の本文 → 無題スライドを開始
            current = {"title": "", "lines": [], "illustration": ""}

        # 見出し4以降・箇条書き・番号付きの記号を除去して本文化
        stripped = re.sub(r"^#{4,}\s*", "", stripped)
        stripped = re.sub(r"^[-*+]\s+", "", stripped)
        stripped = re.sub(r"^\d+[.)]\s+", "", stripped)
        if stripped:
            current["lines"].append(stripped)

    if current is not None:
        slides.append(current)

    return slides


def generate_slide_plan(project_name: str) -> list[dict]:
    """保存済み Markdown から構成案を生成して plan.json に保存する。"""
    md = read_markdown(project_name)
    plan = parse_markdown_to_slides(md)
    plan_path(project_name).write_text(
        json.dumps(plan, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    return plan


def load_plan(project_name: str) -> list[dict]:
    p = plan_path(project_name)
    if not p.exists():
        return []
    try:
        return json.loads(p.read_text(encoding="utf-8"))
    except Exception:
        return []


# ---------------------------------------------------------------------------
# ② 1スライドの画像生成
# ---------------------------------------------------------------------------

def _build_image_prompt(
    project_prompt: str, slide: dict, slide_no: int, total: int, fix: str = "",
    has_reference: bool = False,
) -> str:
    title = slide.get("title", "")
    lines = slide.get("lines", []) or []
    illustration = (slide.get("illustration") or "").strip()

    lines_block = "\n".join(f"・{l}" for l in lines)

    prompt = (
        f"これはプレゼンテーションのスライド（{slide_no}/{total}枚目）です。\n"
        "16:9 のスライド1枚として、見やすく洗練されたレイアウトでデザインしてください。\n"
    )
    if has_reference:
        prompt += (
            "添付した参照画像のデザイン（配色・フォントの雰囲気・レイアウトの方向性）を"
            "踏襲してください。\n"
        )
    prompt += (
        "【余白ルール・重要】見出し・本文・図版など、すべての要素を画像の内側に"
        "余裕をもって収め、上下左右の端に文字を接触させないでください。\n\n"
        f"【スライド見出し（正確に・誤字なく日本語で描画）】\n{title}\n\n"
    )
    if lines_block:
        prompt += f"【本文（正確に・誤字なく日本語で描画。短く要点を）】\n{lines_block}\n\n"
    if illustration:
        prompt += f"【添えるイラスト・図版】\n{illustration}\n（図版内に余計な説明文は入れない）\n\n"
    if project_prompt:
        prompt += f"【デザイン・トーンの指示】\n{project_prompt}\n\n"
    if fix.strip():
        prompt += f"【修正指示（最優先で反映）】\n{fix.strip()}\n"

    return prompt[:30000]


def _fit_to_16x9(img: Image.Image) -> Image.Image:
    """生成画像(3:2)を 16:9 に「縦圧縮でフィット」させる（クロップしない）。

    クロップだと最上端の見出し・最下端の要素のどちらかが切れるため、
    上下を一切切らずに縦方向へ軽く圧縮して 16:9 に収める。
    """
    img = img.convert("RGB")
    return img.resize((TARGET_W, TARGET_H), Image.LANCZOS)


def _sample_background_color(img: Image.Image) -> tuple[int, int, int]:
    """画像の端付近（上辺・左右の内側）から背景色を推定する。"""
    w, h = img.size

    def avg(x0: float, y0: float, x1: float, y1: float) -> tuple[int, int, int]:
        box = (int(x0 * w), int(y0 * h), int(x1 * w), int(y1 * h))
        px = img.crop(box).resize((1, 1), Image.LANCZOS).getpixel((0, 0))
        return px[:3]

    samples = [
        avg(0.10, 0.02, 0.90, 0.05),   # 上辺の内側
        avg(0.015, 0.10, 0.03, 0.90),  # 左辺の内側
        avg(0.97, 0.10, 0.985, 0.90),  # 右辺の内側
    ]
    return tuple(sum(s[i] for s in samples) // len(samples) for i in range(3))


def _apply_negative_space(img: Image.Image, ns: dict) -> Image.Image:
    """右下のネガティブスペースを背景色で軽く均し、きれいな余白に仕上げる。

    外周フレームは inset 分だけ残し、内側の矩形だけをぼかしマスクで合成する。
    """
    img = img.convert("RGB")
    w, h = img.size
    right = float(ns.get("right", 0.15))
    bottom = float(ns.get("bottom", 0.20))
    inset = float(ns.get("inset", 0.014))
    feather = float(ns.get("feather", 0.018))
    color = tuple(ns["color"]) if ns.get("color") else _sample_background_color(img)

    margin = int(inset * w)
    x0 = int((1 - right) * w)
    y0 = int((1 - bottom) * h)
    x1 = w - margin
    y1 = h - margin
    radius = int(min(w, h) * 0.016)

    mask = Image.new("L", (w, h), 0)
    ImageDraw.Draw(mask).rounded_rectangle([x0, y0, x1, y1], radius=radius, fill=255)
    mask = mask.filter(ImageFilter.GaussianBlur(max(1, int(feather * w))))

    fill = Image.new("RGB", (w, h), color)
    return Image.composite(fill, img, mask)


def generate_slide_image(
    project_name: str,
    index: int,
    slide: dict,
    total: int,
    project_prompt: str,
    fix: str = "",
) -> Path:
    """1スライドを GPT Image 2.0 で生成し、16:9 にフィットして保存する。"""
    client = _client()
    ref = reference_path(project_name)
    has_ref = ref.exists()
    prompt = _build_image_prompt(
        project_prompt, slide, index + 1, total, fix=fix, has_reference=has_ref
    )

    if has_ref:
        with open(ref, "rb") as f:
            result = client.images.edit(
                model=SLIDE_IMAGE_MODEL,
                image=[f],
                prompt=prompt,
                size=GEN_SIZE,
            )
    else:
        result = client.images.generate(
            model=SLIDE_IMAGE_MODEL,
            prompt=prompt,
            size=GEN_SIZE,
        )

    b64 = result.data[0].b64_json
    img_bytes = base64.b64decode(b64)

    img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
    img = _fit_to_16x9(img)

    # プロジェクト設定に右下ネガティブスペースがあれば後処理で均す
    ns = read_project_config(project_name).get("negative_space")
    if ns:
        img = _apply_negative_space(img, ns)

    out = slide_png_path(project_name, index)
    img.save(out, "PNG")
    return out


# ---------------------------------------------------------------------------
# ③ PPTX 組立
# ---------------------------------------------------------------------------

def build_pptx(project_name: str, num_slides: int) -> Path:
    """slide_NN.png を 1枚1スライドの 16:9 PPTX に組み立てる。"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for i in range(num_slides):
        png = slide_png_path(project_name, i)
        if not png.exists():
            continue
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(png), 0, 0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    out = pptx_path(project_name)
    prs.save(str(out))
    return out
